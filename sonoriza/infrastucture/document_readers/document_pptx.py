from pathlib import Path
from pptx import Presentation
from sonoriza.core.services.document_reader import IDocumentReader
from sonoriza.core.enums import DocumentType

class PPTXReader(IDocumentReader):
    @property
    def supported_formats(self) -> list[str]:
        return [DocumentType.PPTX]

    def extract_text(self, file_path: Path) -> str:
        """Extrae texto de presentaciones PPTX usando python-pptx"""
        prs = Presentation(file_path)
        full_text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        
        return '\n'.join(full_text)